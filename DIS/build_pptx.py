#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=1.0.0",
#   "beautifulsoup4>=4.12.0",
#   "lxml>=5.0.0",
#   "Pillow>=10.0.0",
# ]
# ///
"""Convert DIS 2026 HTML presentation → PowerPoint with speaker notes."""

from __future__ import annotations
import os
import re
from pathlib import Path
from urllib.parse import unquote

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x08, 0x0c, 0x18)
BG2     = RGBColor(0x11, 0x18, 0x27)
ACCENT  = RGBColor(0x63, 0x66, 0xf1)
ACCENT3 = RGBColor(0x06, 0xb6, 0xd4)
TEXT    = RGBColor(0xf1, 0xf5, 0xf9)
TEXT2   = RGBColor(0x94, 0xa3, 0xb8)
GREEN   = RGBColor(0x10, 0xb9, 0x81)
ORANGE  = RGBColor(0xf5, 0x9e, 0x0b)
BORDER  = RGBColor(0x2a, 0x30, 0x40)

BADGE_MAP: dict[str, tuple[RGBColor, RGBColor]] = {
    "badge-google":  (RGBColor(0x60, 0xa5, 0xfa), RGBColor(0x07, 0x10, 0x28)),
    "badge-meta":    (RGBColor(0x93, 0xc5, 0xfd), RGBColor(0x06, 0x12, 0x2e)),
    "badge-nvidia":  (RGBColor(0x86, 0xef, 0xac), RGBColor(0x04, 0x14, 0x08)),
    "badge-spotify": (RGBColor(0x4a, 0xde, 0x80), RGBColor(0x04, 0x14, 0x08)),
    "badge-misc":    (RGBColor(0xc0, 0x84, 0xfc), RGBColor(0x1e, 0x08, 0x38)),
    "badge-siemens": (RGBColor(0x2d, 0xd4, 0xbf), RGBColor(0x02, 0x16, 0x16)),
    "badge-ica":     (RGBColor(0xfc, 0xa5, 0xa5), RGBColor(0x26, 0x05, 0x05)),
    "keynote":       (RGBColor(0x22, 0xd3, 0xee), RGBColor(0x02, 0x14, 0x1e)),
    "tkwy":          (RGBColor(0xf8, 0x71, 0x71), RGBColor(0x26, 0x05, 0x05)),
    "day1":          (RGBColor(0x10, 0xb9, 0x81), RGBColor(0x04, 0x14, 0x0a)),
    "day2":          (RGBColor(0x81, 0x8c, 0xf8), RGBColor(0x0c, 0x0c, 0x2e)),
    "day3":          (RGBColor(0xfb, 0xbf, 0x24), RGBColor(0x18, 0x10, 0x02)),
}

W  = Inches(13.333)
H  = Inches(7.5)
M  = Inches(0.55)   # slide margin

# ── low-level helpers ─────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height,
         fill=BG2, line=None, lw=Pt(0.5)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape


def txb(slide, text: str, left, top, width, height, *,
        size=13, bold=False, italic=False,
        color=TEXT2, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_badge(slide, text: str, cls: str,
              left=M, top=M, width=Inches(6.0)) -> None:
    fg, bg = BADGE_MAP.get(cls, (ACCENT3, BG2))
    h = Inches(0.28)
    shape = rect(slide, left, top, width, h, fill=bg, line=fg, lw=Pt(0.6))
    tf = shape.text_frame
    tf.margin_left   = Inches(0.08)
    tf.margin_top    = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = text.upper()
    r.font.size      = Pt(7.5)
    r.font.bold      = True
    r.font.color.rgb = fg


def add_notes(slide, text: str) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_image(slide, path: str,
              left, top, width, height) -> bool:
    if not path or not os.path.exists(path):
        if path:
            print(f"  ⚠  not found: {path}")
        return False
    try:
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    except Exception:
        pass
    # MPO / unsupported JPEG variant — re-encode via Pillow
    try:
        import io
        from PIL import Image
        img = Image.open(path)
        if hasattr(img, 'seek'):
            img.seek(0)          # first frame of MPO/animated
        buf = io.BytesIO()
        img.convert('RGB').save(buf, format='JPEG', quality=92)
        buf.seek(0)
        slide.shapes.add_picture(buf, left, top, width, height)
        return True
    except Exception as e:
        print(f"  ⚠  image error {Path(path).name}: {e}")
    return False


def bullet_box(slide, items: list[str],
               left, top, width, height, size=12) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = "→  "
        r1.font.size  = Pt(size - 1)
        r1.font.color.rgb = ACCENT3
        r2 = p.add_run()
        r2.text = item
        r2.font.size  = Pt(size)
        r2.font.color.rgb = TEXT2


# ── HTML extraction helpers ───────────────────────────────────────────────────

def clean(s: str) -> str:
    return re.sub(r'\s+', ' ', s or '').strip()


def get_badge(div: Tag) -> tuple[str, str]:
    for el in div.find_all(True):
        for c in el.get('class', []):
            if c in BADGE_MAP:
                return clean(el.get_text()), c
    return '', 'keynote'


def get_title(div: Tag) -> str:
    for tag in ('h1', 'h2', 'h3'):
        el = div.find(tag)
        if el:
            return clean(el.get_text())
    return ''


def get_lead(div: Tag) -> str:
    el = div.find(class_='lead')
    return clean(el.get_text()) if el else ''


def get_bullets(div: Tag) -> list[str]:
    ul = div.find('ul', class_='bullets')
    if not ul:
        return []
    return [clean(li.get_text()) for li in ul.find_all('li')]


def get_stats(div: Tag) -> list[tuple[str, str]]:
    out = []
    for c in div.find_all(class_='stat-card'):
        num = c.find(class_='stat-number')
        lbl = c.find(class_='stat-label')
        out.append((
            clean(num.get_text()) if num else '',
            clean(lbl.get_text()).replace('\n', ' ') if lbl else '',
        ))
    return out


def get_cards(div: Tag) -> list[tuple[str, str, str]]:
    out = []
    for c in div.find_all(class_='card'):
        if 'stat-card' in c.get('class', []):
            continue
        icon = c.find(class_='icon')
        h4   = c.find('h4')
        p    = c.find('p')
        out.append((
            clean(icon.get_text()) if icon else '',
            clean(h4.get_text())   if h4   else '',
            clean(p.get_text())    if p    else '',
        ))
    return out


def get_pipeline(div: Tag) -> list[tuple[str, str, str]]:
    out = []
    for s in div.find_all(class_='pipeline-step'):
        num  = s.find(class_='step-num')
        name = s.find(class_='step-name')
        desc = s.find(class_='step-desc')
        out.append((
            clean(num.get_text())  if num  else '',
            clean(name.get_text()) if name else '',
            clean(desc.get_text()) if desc else '',
        ))
    return out


def get_day_cards(div: Tag) -> list[tuple[str, str]]:
    out = []
    for c in div.find_all(class_='day-card'):
        h4 = c.find('h4')
        p  = c.find('p')
        out.append((
            clean(h4.get_text()) if h4 else '',
            clean(p.get_text())  if p  else '',
        ))
    return out


def get_tkwy(div: Tag) -> list[tuple[str, str, str]]:
    out = []
    for it in div.find_all(class_='tkwy-item'):
        num = it.find(class_='tkwy-num')
        h4  = it.find('h4')
        p   = it.find('p')
        out.append((
            clean(num.get_text()) if num else '',
            clean(h4.get_text())  if h4  else '',
            clean(p.get_text())   if p   else '',
        ))
    return out


def get_image(div: Tag, base: Path) -> str | None:
    for img in div.find_all('img'):
        src = img.get('src', '')
        if not src:
            continue
        for candidate in (base / src, base / unquote(src)):
            if candidate.exists():
                return str(candidate)
    return None


def get_quote(div: Tag) -> str:
    q = div.find('blockquote')
    return clean(q.get_text()) if q else ''


# ── slide builders ────────────────────────────────────────────────────────────

def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_title_slide(prs, div, base):
    slide = blank(prs)
    set_bg(slide)
    add_notes(slide, div.get('data-notes', ''))

    badge, cls = get_badge(div)
    if badge:
        add_badge(slide, badge, cls,
                  left=Inches(3.5), top=Inches(0.9), width=Inches(6.3))

    title = get_title(div)
    if title:
        txb(slide, title, M, Inches(1.35), W - 2*M, Inches(1.6),
            size=52, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    sub = div.find(class_='subtitle')
    if sub:
        txb(slide, clean(sub.get_text()), M, Inches(3.05), W - 2*M, Inches(0.5),
            size=16, color=TEXT2, align=PP_ALIGN.CENTER)

    meta = div.find(class_='meta-date')
    if meta:
        txb(slide, clean(meta.get_text()), M, Inches(3.6), W - 2*M, Inches(0.35),
            size=12, color=TEXT2, align=PP_ALIGN.CENTER)

    stats = get_stats(div)
    if stats:
        n   = len(stats)
        cw  = Inches(1.9)
        gap = Inches(0.22)
        tw  = n * cw + (n - 1) * gap
        x0  = (W - tw) / 2
        y0  = Inches(4.15)
        for i, (num, lbl) in enumerate(stats):
            x = x0 + i * (cw + gap)
            rect(slide, x, y0, cw, Inches(1.25), fill=BG2, line=BORDER)
            txb(slide, num, x, y0 + Inches(0.08), cw, Inches(0.62),
                size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
            txb(slide, lbl, x, y0 + Inches(0.72), cw, Inches(0.45),
                size=9, color=TEXT2, align=PP_ALIGN.CENTER)


def build_quote_slide(prs, div, base):
    slide = blank(prs)
    set_bg(slide)
    add_notes(slide, div.get('data-notes', ''))

    badge, cls = get_badge(div)
    if badge:
        add_badge(slide, badge, cls,
                  left=Inches(3.5), top=M, width=Inches(6.3))

    q = get_quote(div)
    if q:
        txb(slide, f'"{q}"', M, Inches(2.0), W - 2*M, Inches(2.5),
            size=30, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    for i, attr in enumerate(div.find_all(class_='quote-attr')):
        txb(slide, clean(attr.get_text()),
            M, Inches(4.6) + Inches(0.45 * i), W - 2*M, Inches(0.42),
            size=14, color=TEXT2, align=PP_ALIGN.CENTER)


def build_two_col_slide(prs, div, base):
    slide = blank(prs)
    set_bg(slide)
    add_notes(slide, div.get('data-notes', ''))

    left  = div.find(class_='col-left') or div
    right = div.find(class_='col-right')

    # right panel background
    rx = Inches(6.85)
    rw = W - rx
    rect(slide, rx - Inches(0.05), 0, rw + Inches(0.05), H, fill=BG2)

    col_w = rx - M - Inches(0.25)
    y = M

    badge, cls = get_badge(left)
    if badge:
        add_badge(slide, badge, cls, left=M, top=y, width=Inches(5.9))
        y += Inches(0.37)

    title = get_title(left)
    if title:
        txb(slide, title, M, y, col_w, Inches(1.05),
            size=21, bold=True, color=TEXT)
        y += Inches(1.1)

    lead = get_lead(left)
    if lead:
        txb(slide, lead, M, y, col_w, Inches(0.38),
            size=12, color=TEXT2)
        y += Inches(0.42)

    # Stats row
    stats = get_stats(left)
    if stats:
        n  = len(stats)
        sw = (col_w - Inches(0.1) * (n - 1)) / n
        for i, (num, lbl) in enumerate(stats):
            sx = M + i * (sw + Inches(0.1))
            rect(slide, sx, y, sw, Inches(0.85), fill=BG2, line=BORDER)
            txb(slide, num, sx, y, sw, Inches(0.5),
                size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
            txb(slide, lbl, sx, y + Inches(0.5), sw, Inches(0.32),
                size=8, color=TEXT2, align=PP_ALIGN.CENTER)
        y += Inches(1.0)

    # Pipeline
    pipeline = get_pipeline(left)
    if pipeline:
        n  = len(pipeline)
        pw = min((col_w - Inches(0.32) * (n - 1)) / n, Inches(1.25))
        for i, (pnum, pname, pdesc) in enumerate(pipeline):
            px = M + i * (pw + Inches(0.32))
            if px + pw > M + col_w:
                break
            rect(slide, px, y, pw, Inches(0.78), fill=BG2,
                 line=RGBColor(0x3a, 0x3f, 0x80))
            txb(slide, pnum, px, y, pw, Inches(0.22),
                size=7, color=ACCENT, align=PP_ALIGN.CENTER)
            txb(slide, pname, px, y + Inches(0.22), pw, Inches(0.32),
                size=9, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
            if pdesc:
                txb(slide, pdesc, px, y + Inches(0.54), pw, Inches(0.22),
                    size=7, color=TEXT2, align=PP_ALIGN.CENTER)
            if i < n - 1:
                txb(slide, '→', px + pw + Inches(0.02), y + Inches(0.24),
                    Inches(0.28), Inches(0.28),
                    size=13, color=ACCENT3, align=PP_ALIGN.CENTER)
        y += Inches(0.92)

    bullets = get_bullets(left)
    cards   = get_cards(left)

    if bullets:
        bullet_box(slide, bullets, M, y, col_w, H - y - Inches(0.15), size=12)

    elif cards:
        n     = len(cards)
        cols  = min(n, 2)
        rows  = (n + cols - 1) // cols
        avail = H - y - Inches(0.12)
        cw    = (col_w - Inches(0.12) * (cols - 1)) / cols
        ch    = (avail  - Inches(0.12) * (rows - 1)) / rows
        for i, (icon, ctitle, cdesc) in enumerate(cards):
            cx = M + (i % cols) * (cw + Inches(0.12))
            cy = y + (i // cols) * (ch + Inches(0.12))
            rect(slide, cx, cy, cw, ch, fill=BG2, line=BORDER)
            txb(slide, f"{icon}  {ctitle}",
                cx + Inches(0.1), cy + Inches(0.07), cw - Inches(0.18), Inches(0.28),
                size=10, bold=True, color=TEXT)
            txb(slide, cdesc,
                cx + Inches(0.1), cy + Inches(0.38), cw - Inches(0.18), ch - Inches(0.5),
                size=9, color=TEXT2)

    # Right image
    img = get_image(right or div, base)
    if img:
        ip = Inches(0.22)
        add_image(slide, img, rx + ip, ip, rw - 2 * ip - M * 0.5, H - 2 * ip)


def build_regular_slide(prs, div, base):
    slide = blank(prs)
    set_bg(slide)
    add_notes(slide, div.get('data-notes', ''))

    badge, cls = get_badge(div)
    y = M

    if badge:
        add_badge(slide, badge, cls, left=M, top=y, width=Inches(6.5))
        y += Inches(0.37)

    title = get_title(div)
    if title:
        txb(slide, title, M, y, W - 2 * M, Inches(0.92),
            size=28, bold=True, color=TEXT)
        y += Inches(0.98)

    # ── summit overview (day cards) ──
    day_cards = get_day_cards(div)
    if day_cards:
        n  = len(day_cards)
        cw = (W - 2 * M - Inches(0.25) * (n - 1)) / n
        ch = H - y - Inches(0.2)
        day_colors = [GREEN, RGBColor(0x81, 0x8c, 0xf8), ORANGE]
        for i, (dh, dp) in enumerate(day_cards):
            cx = M + i * (cw + Inches(0.25))
            color = day_colors[i % len(day_colors)]
            rect(slide, cx, y, cw, ch, fill=BG2, line=color)
            txb(slide, dh, cx + Inches(0.15), y + Inches(0.12),
                cw - Inches(0.3), Inches(0.32),
                size=10, bold=True, color=color)
            txb(slide, dp, cx + Inches(0.15), y + Inches(0.48),
                cw - Inches(0.3), ch - Inches(0.62),
                size=11, color=TEXT2)
        return

    # ── stats ──
    stats = get_stats(div)
    if stats:
        n   = len(stats)
        sw  = min((W - 2 * M - Inches(0.18) * (n - 1)) / n, Inches(2.5))
        tw  = n * sw + (n - 1) * Inches(0.18)
        x0  = (W - tw) / 2
        for i, (num, lbl) in enumerate(stats):
            sx = x0 + i * (sw + Inches(0.18))
            rect(slide, sx, y, sw, Inches(1.1), fill=BG2, line=BORDER)
            txb(slide, num, sx, y + Inches(0.05), sw, Inches(0.62),
                size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
            txb(slide, lbl, sx, y + Inches(0.68), sw, Inches(0.38),
                size=9, color=TEXT2, align=PP_ALIGN.CENTER)
        y += Inches(1.25)

    # ── bullets ──
    bullets = get_bullets(div)
    if bullets:
        bullet_box(slide, bullets, M, y, W - 2 * M, H - y - Inches(0.2), size=14)
        return

    # ── feature cards (e.g. Siemens, Thank You) ──
    cards = get_cards(div)
    if cards:
        n    = len(cards)
        cols = 4 if n > 4 else min(n, 3)
        rows = (n + cols - 1) // cols
        avail_h = H - y - Inches(0.15)
        cw = (W - 2 * M - Inches(0.15) * (cols - 1)) / cols
        ch = (avail_h - Inches(0.15) * (rows - 1)) / rows
        for i, (icon, ctitle, cdesc) in enumerate(cards):
            cx = M + (i % cols) * (cw + Inches(0.15))
            cy = y + (i // cols) * (ch + Inches(0.15))
            rect(slide, cx, cy, cw, ch, fill=BG2, line=BORDER)
            txb(slide, f"{icon}  {ctitle}",
                cx + Inches(0.12), cy + Inches(0.1), cw - Inches(0.22), Inches(0.32),
                size=11, bold=True, color=TEXT)
            txb(slide, cdesc,
                cx + Inches(0.12), cy + Inches(0.46), cw - Inches(0.22), ch - Inches(0.58),
                size=10, color=TEXT2)
        return

    # ── numbered takeaways ──
    tkwy = get_tkwy(div)
    if tkwy:
        n     = len(tkwy)
        avail = H - y - Inches(0.15)
        ih    = avail / n - Inches(0.08)
        for i, (num, th, tp) in enumerate(tkwy):
            iy = y + i * (ih + Inches(0.08))
            rect(slide, M, iy, W - 2 * M, ih, fill=BG2, line=ACCENT)
            txb(slide, num, M + Inches(0.1), iy + Inches(0.02), Inches(0.45), ih,
                size=16, bold=True, color=ACCENT)
            txb(slide, th,
                M + Inches(0.6), iy + Inches(0.02), W - 2 * M - Inches(0.7), Inches(0.3),
                size=11, bold=True, color=TEXT)
            txb(slide, tp,
                M + Inches(0.6), iy + Inches(0.34), W - 2 * M - Inches(0.7), ih - Inches(0.38),
                size=9.5, color=TEXT2)

    # ── quote (fallback) ──
    q = get_quote(div)
    if q and not title:
        txb(slide, f'"{q}"', M, Inches(2.2), W - 2 * M, Inches(2.0),
            size=28, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def build_mantra_slide(prs, div, base):
    slide = blank(prs)
    set_bg(slide)
    add_notes(slide, div.get('data-notes', ''))

    badge, cls = get_badge(div)
    y = M
    if badge:
        add_badge(slide, badge, cls, left=M, top=y, width=Inches(6.5))
        y += Inches(0.37)

    title = get_title(div)
    if title:
        txb(slide, title, M, y, W - 2 * M, Inches(0.8),
            size=26, bold=True, color=TEXT)
        y += Inches(0.88)

    img = get_image(div, base)
    if img:
        add_image(slide, img, M, y, W - 2 * M, Inches(2.1))
        y += Inches(2.2)

    mantras = [
        ("ABT\nAlways Be Testing", "Small iterations · Fail fast · Learn", ACCENT),
        ('"What you cannot\nmeasure doesn\'t exist"', "Track what matters · Instrument everything",
         RGBColor(0x81, 0x8c, 0xf8)),
        ("Progress\n> Perfection", "Ship · Observe · Improve\nDon't wait for perfect", ORANGE),
    ]
    n   = len(mantras)
    sw  = (W - 2 * M - Inches(0.2) * (n - 1)) / n
    sh  = H - y - Inches(0.15)
    for i, (label, desc, color) in enumerate(mantras):
        sx = M + i * (sw + Inches(0.2))
        rect(slide, sx, y, sw, sh, fill=BG2, line=BORDER)
        txb(slide, label, sx, y + Inches(0.1), sw, Inches(0.72),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        txb(slide, desc, sx, y + Inches(0.85), sw, sh - Inches(0.95),
            size=10, color=TEXT2, align=PP_ALIGN.CENTER)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    base      = Path(__file__).parent
    html_path = base / "presentation.html"
    out_path  = base / "DIS_2026_Key_Takeaways.pptx"

    print(f"Reading {html_path.name} …")
    soup      = BeautifulSoup(html_path.read_text(encoding="utf-8"), "lxml")
    slide_divs = soup.find_all("div", class_="slide")
    print(f"Found {len(slide_divs)} slides\n")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for i, div in enumerate(slide_divs, 1):
        classes    = div.get("class", [])
        title_text = get_title(div) or "(no title)"

        if "title-slide" in classes:
            kind = "title"
        elif "quote-slide" in classes:
            kind = "quote"
        elif "two-col" in classes:
            kind = "two-col"
        elif "Three Mantras" in title_text:
            kind = "mantra"
        else:
            kind = "regular"

        print(f"  [{i:02d}] {kind:<8}  {title_text[:60]}")

        if kind == "title":
            build_title_slide(prs, div, base)
        elif kind == "quote":
            build_quote_slide(prs, div, base)
        elif kind == "two-col":
            build_two_col_slide(prs, div, base)
        elif kind == "mantra":
            build_mantra_slide(prs, div, base)
        else:
            build_regular_slide(prs, div, base)

    prs.save(str(out_path))
    print(f"\n✓  Saved → {out_path.name}")


if __name__ == "__main__":
    main()
